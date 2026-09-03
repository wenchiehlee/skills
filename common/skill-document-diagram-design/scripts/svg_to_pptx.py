#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
svg_to_pptx.py — 把一組已轉出的 PNG 圖表組成 PowerPoint（.pptx）

輸入的圖片必須已經是 PNG（原生管線用 svg_to_png.py、PlantUML 管線用
render_plantuml.py fetch --fmt png 轉出），本腳本不做任何向量轉點陣的工作。

版型固定：16:9、每筆一張投影片、標題置頂、圖片等比縮放置中、可選講者備忘稿。
需要客製版型時，直接把 PNG 貼進既有 PPTX 範本更實際，不在此腳本擴充。

用法：
  python scripts/svg_to_pptx.py manifest.json --out deck.pptx

manifest.json 格式：
  [
    {"title": "投資十步流程", "image": "flowchart.png", "notes": "對應 SKILL.md 十步判斷"},
    {"title": "Investment Stack Vision", "image": "InvestmentStackVision.png"}
  ]
image 路徑為相對於 manifest.json 所在目錄的路徑。
"""
from __future__ import annotations

import argparse
import json
import platform
import sys
from pathlib import Path

if platform.system() == "Windows":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass


def build_pptx(manifest_path: Path, out_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("[svg_to_pptx] 缺少 python-pptx，請先執行：pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    from PIL import Image

    if not manifest_path.exists():
        print(f"[svg_to_pptx] 找不到 manifest：{manifest_path}", file=sys.stderr)
        sys.exit(1)

    entries = json.loads(manifest_path.read_text(encoding="utf-8"))
    base_dir = manifest_path.parent

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    title_h = Inches(0.9)
    margin = Inches(0.5)

    for entry in entries:
        title = entry.get("title", "")
        image_rel = entry.get("image")
        notes = entry.get("notes")

        if not image_rel:
            print(f"[svg_to_pptx] 略過缺少 image 欄位的項目：{entry}", file=sys.stderr)
            continue

        image_path = base_dir / image_rel
        if not image_path.exists():
            print(f"[svg_to_pptx] 找不到圖片，略過：{image_path}", file=sys.stderr)
            continue

        slide = prs.slides.add_slide(blank_layout)

        if title:
            tb = slide.shapes.add_textbox(margin, Inches(0.2), slide_w - 2 * margin, title_h)
            tf = tb.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.bold = True

        with Image.open(image_path) as img:
            img_w_px, img_h_px = img.size

        area_top = title_h + Inches(0.3) if title else margin
        area_w = slide_w - 2 * margin
        area_h = slide_h - area_top - margin

        img_ratio = img_w_px / img_h_px
        area_ratio = area_w / area_h

        if img_ratio > area_ratio:
            draw_w = area_w
            draw_h = int(area_w / img_ratio)
        else:
            draw_h = area_h
            draw_w = int(area_h * img_ratio)

        left = margin + (area_w - draw_w) // 2
        top = area_top + (area_h - draw_h) // 2

        slide.shapes.add_picture(str(image_path), left, top, width=draw_w, height=draw_h)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"[svg_to_pptx] 已輸出：{out_path}（{len(entries)} 張投影片）")


def main() -> None:
    parser = argparse.ArgumentParser(description="把一組 PNG 圖表組成 PowerPoint")
    parser.add_argument("manifest", help="manifest.json 路徑")
    parser.add_argument("--out", required=True, help="輸出 .pptx 路徑")
    args = parser.parse_args()

    build_pptx(Path(args.manifest), Path(args.out))


if __name__ == "__main__":
    main()
