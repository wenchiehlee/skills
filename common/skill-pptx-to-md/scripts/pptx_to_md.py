#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx_to_md.py — 將 PowerPoint (.pptx) 簡報轉換為 Markdown 格式

使用 python-pptx 解析投影片中的標題、內文、項目符號、表格、
講者備忘稿，並可選擇將投影片內嵌圖片抽取到本地資料夾，
於 Markdown 中以相對路徑連結。

使用方式：

    pip install python-pptx

    # 輸出到 stdout
    python scripts/pptx_to_md.py path/to/deck.pptx > output.md

    # 同時抽取圖片到 images/ 資料夾（Markdown 內自動加上圖片連結）
    python scripts/pptx_to_md.py path/to/deck.pptx --extract-images images > output.md

作為模組使用：

    from scripts.pptx_to_md import convert_pptx_to_markdown
    markdown_text = convert_pptx_to_markdown("deck.pptx", image_dir="images")
"""
import argparse
import platform
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Fix Windows console encoding for Chinese characters
if platform.system() == "Windows":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass

SLIDE_MARKER = "<!-- SLIDE:{slide} -->"


def _paragraph_to_md(paragraph) -> str:
    text = "".join(run.text for run in paragraph.runs) or paragraph.text
    text = text.strip()
    if not text:
        return ""
    level = getattr(paragraph, "level", 0) or 0
    indent = "  " * level
    return f"{indent}- {text}"


def _text_frame_to_md(shape, is_title: bool) -> list:
    lines = []
    tf = shape.text_frame
    paragraphs = [p for p in tf.paragraphs if p.text.strip()]
    if is_title:
        title_text = tf.text.strip()
        if title_text:
            lines.append(f"## {title_text}")
            lines.append("")
        return lines

    for para in paragraphs:
        md_line = _paragraph_to_md(para)
        if md_line:
            lines.append(md_line)
    if lines:
        lines.append("")
    return lines


def _table_to_md(shape) -> list:
    table = shape.table
    lines = []
    rows = list(table.rows)
    if not rows:
        return lines

    def _row_text(row):
        return [cell.text.strip().replace("\n", "<br>") for cell in row.cells]

    header = _row_text(rows[0])
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(_row_text(row)) + " |")
    lines.append("")
    return lines


def _extract_image(shape, image_dir: Path, slide_idx: int, img_idx: int) -> str:
    image = shape.image
    ext = image.ext or "png"
    filename = f"slide{slide_idx:03d}_img{img_idx:02d}.{ext}"
    dest = image_dir / filename
    image_dir.mkdir(parents=True, exist_ok=True)
    with open(dest, "wb") as f:
        f.write(image.blob)
    return filename


def _shape_to_md(shape, slide_idx: int, img_counter: list, image_dir: Path = None) -> list:
    lines = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            lines.extend(_shape_to_md(sub_shape, slide_idx, img_counter, image_dir))
        return lines

    if shape.has_table:
        lines.extend(_table_to_md(shape))
        return lines

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and image_dir is not None:
        img_counter[0] += 1
        filename = _extract_image(shape, image_dir, slide_idx, img_counter[0])
        rel_path = f"{image_dir.name}/{filename}"
        lines.append(f"![slide{slide_idx} image]({rel_path})")
        lines.append("")
        return lines

    if shape.has_text_frame and shape.text_frame.text.strip():
        # 非標題文字方塊一律以項目符號輸出
        lines.extend(_text_frame_to_md(shape, is_title=False))

    return lines


def convert_pptx_to_markdown(pptx_path, image_dir: str = None, include_notes: bool = True) -> str:
    """
    將 .pptx 簡報轉換為 Markdown 文本。

    :param pptx_path: PPTX 檔案路徑
    :param image_dir: 若指定，投影片內嵌圖片會被抽取到此資料夾（相對於輸出檔案）
    :param include_notes: 是否包含講者備忘稿
    :return: Markdown 文本（含 SLIDE 標記）
    """
    path_obj = Path(pptx_path)
    if not path_obj.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")

    prs = Presentation(str(path_obj))
    source = path_obj.name
    image_out_dir = Path(image_dir) if image_dir else None

    parts = [f"# {path_obj.stem}", ""]

    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(SLIDE_MARKER.format(slide=idx))

        title_shape = slide.shapes.title
        title_text = title_shape.text_frame.text.strip() if title_shape and title_shape.has_text_frame else ""
        parts.append(f"## Slide {idx}" + (f": {title_text}" if title_text else ""))
        parts.append("")

        img_counter = [0]
        title_id = title_shape.shape_id if title_shape is not None else None
        for shape in slide.shapes:
            if title_id is not None and shape.shape_id == title_id:
                continue
            body_lines = _shape_to_md(shape, idx, img_counter, image_out_dir)
            parts.extend(body_lines)

        if include_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append("> **講者備忘稿：**")
                for line in notes_text.splitlines():
                    parts.append(f"> {line}")
                parts.append("")

        parts.append("---")
        parts.append("")

    slide_count = len(prs.slides)
    print(f"[pptx_to_md] {source}：共 {slide_count} 張投影片轉換完成", file=sys.stderr)

    return "\n".join(parts)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="將 PowerPoint (.pptx) 轉換為 Markdown")
    parser.add_argument("pptx", help="PPTX 檔案路徑")
    parser.add_argument("--extract-images", dest="image_dir", default=None,
                        help="抽取投影片內嵌圖片到指定資料夾（相對路徑會寫入 Markdown 圖片連結）")
    parser.add_argument("--no-notes", dest="include_notes", action="store_false",
                        help="不包含講者備忘稿")
    args = parser.parse_args()

    try:
        print(convert_pptx_to_markdown(args.pptx, image_dir=args.image_dir, include_notes=args.include_notes))
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
